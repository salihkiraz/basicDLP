import os
import re
import docx
import openpyxl
from pptx import Presentation
from pdfminer.high_level import extract_text as extract_text_from_pdf
from colorama import Fore, Style
from tqdm import tqdm
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext
from tkinter import ttk  

def load_name_list(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        name_list = [line.strip() for line in file.readlines()]
    return name_list

def extract_text_from_docx(docx_path):
    try:
        doc = docx.Document(docx_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except docx.opc.exceptions.PackageNotFoundError:
        return ""

def extract_text_from_pdf_path(pdf_path):
    try:
        return extract_text_from_pdf(pdf_path)
    except Exception as e:
        return ""

def extract_text_from_excel(excel_path):
    try:
        wb = openpyxl.load_workbook(excel_path)
        full_text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        full_text.append(str(cell.value))
        return '\n'.join(full_text)
    except Exception as e:
        return ""

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text)
    except Exception as e:
        return ""

def find_matches(text, regex_patterns, name_list):
    matches = {}
    for label, regex in regex_patterns.items():
        found = re.findall(regex, text)
        if found:
            matches[label] = found

    # Ad Soyad araması
    ad_soyad_matches = []
    for ad_soyad in name_list:
        if ad_soyad in text:
            ad_soyad_matches.append(ad_soyad)

    if ad_soyad_matches:
        matches["Ad Soyad"] = ad_soyad_matches

    return matches

def extract_text_from_txt(txt_path):
    try:
        with open(txt_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        print(f"Error: Could not extract text from TXT file: {txt_path}. Error: {e}")
        return ""

def extract_text_from_file(file_path):
    if file_path.endswith('.docx') or file_path.endswith('.doc'):
        return extract_text_from_docx(file_path)
    elif file_path.endswith('.pdf'):
        return extract_text_from_pdf_path(file_path)
    elif file_path.endswith('.xlsx') or file_path.endswith('.xls'):
        return extract_text_from_excel(file_path)
    elif file_path.endswith('.pptx') or file_path.endswith('.ppt'):
        return extract_text_from_pptx(file_path)
    elif file_path.endswith('.txt'):
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_path}")

def find_files_with_regex(root_dir, regex_patterns, name_list, progress_callback=None):
    matched_files = []

    valid_extensions = {'.docx', '.doc', '.pdf', '.xlsx', '.xls', '.pptx', '.ppt', '.txt'}
    
    file_paths = []
    for dirpath, _, filenames in os.walk(root_dir):
        for filename in filenames:
            if filename.startswith('~$'):  # Skip temporary files
                continue
            file_ext = os.path.splitext(filename)[1].lower()
            if file_ext in valid_extensions:
                file_paths.append(os.path.join(dirpath, filename))

    for i, file_path in enumerate(tqdm(file_paths, desc="Scanning files")):
        if progress_callback:
            progress_callback(i+1, len(file_paths))
        try:
            text = extract_text_from_file(file_path)
            if text:  # Only process if text extraction was successful
                matches = find_matches(text, regex_patterns, name_list)
                if matches:
                    matched_files.append((file_path, matches))
        except (IOError, ValueError) as e:
            print(f"Error: An error occurred while processing {file_path}. Error: {e}")
            continue

    return matched_files

def run_scan(root_directory, name_list_file, output_text_widget, progress_var):
    name_list = load_name_list(name_list_file)

    regex_patterns = {
        "E-mail": r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}',
        "TCKN": r'\b[1-9][0-9]{10}\b',  # VAT NUMBER for TURKEY VAT 11-digit numbers that do not start with 0
        "Telephone": r'\b(\+90\s?)?(5\d{2}\s?\d{3}\s?\d{2}\s?\d{2}|5\d{2}-\d{3}-\d{2}-\d{2}|(\(5\d{2}\)\s?\d{3}\s?\d{2}\s?\d{2})|(\(5\d{2}\)-\d{3}-\d{2}-\d{2}))\b',  #format for Turkey Phone
        "Student No": r'\b\d{10}\b'  # 10-digit school numbers
    }

    def update_progress(current, total):
        progress_var.set(int((current / total) * 100))
        root.update_idletasks()

    matched_files = find_files_with_regex(root_directory, regex_patterns, name_list, update_progress)

    output_text_widget.delete(1.0, tk.END)
    if matched_files:
        output_text_widget.insert(tk.END, f"Files that fit the specified patterns:\n")
        for file_path, matches in matched_files:
            output_text_widget.insert(tk.END, f"\File: {file_path}\n")
            for label, found_matches in matches.items():
                output_text_widget.insert(tk.END, f"  {label}:\n")
                for match in found_matches:
                    output_text_widget.insert(tk.END, f"    {match}\n")
    else:
        output_text_widget.insert(tk.END, "No content matching the specified patterns was found in any file.")

def select_directory():
    directory = filedialog.askdirectory()
    if directory:
        root_directory.set(directory)

def select_name_list_file():
    file_path = filedialog.askopenfilename(filetypes=[("Text files", "*.txt")])
    if file_path:
        name_list_file.set(file_path)

def add_developer_info():
    info_text = """\
    This application is developed using Python and the tkinter library.
    Developer: Salih KİRAZ
    Email: salihk06@gmail.com
    """
    label.config(text=info_text)

root = tk.Tk()
root.title("Basic DLP ")

root_directory = tk.StringVar()
name_list_file = tk.StringVar()
progress_var = tk.IntVar()

frame = tk.Frame(root)
frame.pack(pady=10, padx=10, fill=tk.X)

tk.Label(frame, text="Scan Directory:").grid(row=0, column=0, sticky=tk.W)
tk.Entry(frame, textvariable=root_directory, width=50).grid(row=0, column=1, padx=5, sticky=tk.W)
tk.Button(frame, text="Browse", command=select_directory).grid(row=0, column=2)


tk.Button(root, text="Start Scan", command=lambda: run_scan(root_directory.get(), "name_list.txt", output_text, progress_var)).pack(pady=5)

output_text = scrolledtext.ScrolledText(root, width=100, height=50, wrap=tk.WORD, font=("Arial", 10))
output_text.pack(fill="both", expand=True)


progress_bar = ttk.Progressbar(root, variable=progress_var, maximum=100)
progress_bar.pack(pady=5, padx=10, fill=tk.X)

label = tk.Label(root, width=50, height=10, wraplength=400, justify="left", font=("Arial", 12))
label.pack(fill="both", expand=True)

# Call function to add developer info
add_developer_info()

root.mainloop()
